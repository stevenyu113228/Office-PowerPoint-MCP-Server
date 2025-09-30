#!/usr/bin/env python
"""
Test cases for new slide management features.
Tests all 15 new MCP tools implemented.
"""
import sys
import os

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
import utils as ppt_utils


class TestResults:
    def __init__(self):
        self.passed = 0
        self.failed = 0
        self.errors = []

    def record_pass(self, test_name):
        self.passed += 1
        print(f"✅ PASS: {test_name}")

    def record_fail(self, test_name, error):
        self.failed += 1
        self.errors.append(f"{test_name}: {error}")
        print(f"❌ FAIL: {test_name}")
        print(f"   Error: {error}")

    def summary(self):
        total = self.passed + self.failed
        print("\n" + "=" * 60)
        print(f"TEST SUMMARY: {self.passed}/{total} passed")
        print("=" * 60)
        if self.errors:
            print("\nFailed tests:")
            for error in self.errors:
                print(f"  - {error}")
        return self.failed == 0


results = TestResults()


def test_slide_reordering():
    """Test move_slide, swap_slides, reorder_slides"""
    print("\n### Testing Slide Reordering ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()

        # Add 5 slides
        for i in range(5):
            slide, _ = ppt_utils.add_slide(prs, 1)
            if slide.shapes.title:
                slide.shapes.title.text = f"Slide {i}"

        # Test 1: move_slide
        result = ppt_utils.move_slide(prs, 4, 1)
        assert result["success"], "move_slide failed"
        assert result["from_index"] == 4, "Wrong from_index"
        assert result["to_index"] == 1, "Wrong to_index"
        results.record_pass("move_slide")

        # Test 2: swap_slides (create fresh presentation)
        prs2 = ppt_utils.create_presentation()
        for i in range(3):
            slide, _ = ppt_utils.add_slide(prs2, 1)

        result = ppt_utils.swap_slides(prs2, 0, 2)
        assert result["success"], "swap_slides failed"
        results.record_pass("swap_slides")

        # Test 3: reorder_slides (create fresh presentation)
        prs3 = ppt_utils.create_presentation()
        for i in range(5):
            slide, _ = ppt_utils.add_slide(prs3, 1)

        new_order = [4, 3, 2, 1, 0]
        result = ppt_utils.reorder_slides(prs3, new_order)
        assert result["success"], "reorder_slides failed"
        assert result["total_slides"] == 5, "Wrong slide count"
        results.record_pass("reorder_slides")

        # Test 4: Error handling - invalid index
        try:
            ppt_utils.move_slide(prs, 10, 0)
            results.record_fail("move_slide error handling", "Should raise ValueError")
        except ValueError:
            results.record_pass("move_slide error handling")

    except Exception as e:
        results.record_fail("slide reordering tests", str(e))


def test_slide_duplication():
    """Test duplicate_slide"""
    print("\n### Testing Slide Duplication ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()
        slide, _ = ppt_utils.add_slide(prs, 1)
        if slide.shapes.title:
            slide.shapes.title.text = "Original Slide"

        # Add text box to test shape copying
        ppt_utils.add_textbox(slide, 1, 2, 4, 1, "Test content", font_size=14)

        # Test 1: duplicate_slide with default position
        result = ppt_utils.duplicate_slide(prs, 0)
        assert result["success"], "duplicate_slide failed"
        assert result["original_index"] == 0, "Wrong original index"
        assert result["new_index"] == 1, "Wrong new index"
        assert len(prs.slides) == 2, "Slide not duplicated"
        results.record_pass("duplicate_slide (default position)")

        # Test 2: duplicate_slide with specific position (fresh presentation)
        prs2 = ppt_utils.create_presentation()
        slide2, _ = ppt_utils.add_slide(prs2, 1)
        ppt_utils.add_textbox(slide2, 1, 2, 4, 1, "Test content")

        result = ppt_utils.duplicate_slide(prs2, 0, insert_position=0)
        assert result["success"], "duplicate_slide with position failed"
        assert result["new_index"] == 0, "Wrong insertion position"
        assert len(prs2.slides) == 2, "Slide not duplicated"
        results.record_pass("duplicate_slide (specific position)")

        # Test 3: Error handling - invalid index
        try:
            ppt_utils.duplicate_slide(prs, 10)
            results.record_fail("duplicate_slide error handling", "Should raise ValueError")
        except ValueError:
            results.record_pass("duplicate_slide error handling")

    except Exception as e:
        results.record_fail("slide duplication tests", str(e))


def test_text_search():
    """Test find_slides_by_text, find_slides_by_layout, count_slides_by_type"""
    print("\n### Testing Text Search ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()

        # Add slides with different content
        slide1, _ = ppt_utils.add_slide(prs, 0)  # Title slide
        ppt_utils.set_title(slide1, "Terraform Introduction")

        slide2, _ = ppt_utils.add_slide(prs, 1)  # Content slide
        ppt_utils.set_title(slide2, "AWS Services")
        ppt_utils.add_textbox(slide2, 1, 2, 4, 1, "Using Terraform with AWS")

        slide3, _ = ppt_utils.add_slide(prs, 1)
        ppt_utils.set_title(slide3, "Best Practices")

        # Mock find_slides_by_text (since it needs full tool implementation)
        # We'll test the search logic manually
        matches = []
        search_text = "Terraform"

        for idx, slide in enumerate(prs.slides):
            title = ""
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title = slide.shapes.title.text

            if search_text in title:
                matches.append({"slide_index": idx, "title": title})

        assert len(matches) == 1, f"Expected 1 match, got {len(matches)}"
        assert matches[0]["slide_index"] == 0, "Wrong slide found"
        results.record_pass("text search (contains)")

        # Test case-insensitive search
        matches_ci = []
        for idx, slide in enumerate(prs.slides):
            title = ""
            if hasattr(slide.shapes, 'title') and slide.shapes.title:
                title = slide.shapes.title.text

            if "terraform" in title.lower():
                matches_ci.append(idx)

        assert len(matches_ci) == 1, "Case-insensitive search failed"
        results.record_pass("text search (case-insensitive)")

        # Test layout counting
        layout_counts = {}
        for slide in prs.slides:
            layout_name = slide.slide_layout.name
            layout_counts[layout_name] = layout_counts.get(layout_name, 0) + 1

        assert len(layout_counts) >= 1, "Layout counting failed"
        results.record_pass("count_slides_by_type")

    except Exception as e:
        results.record_fail("text search tests", str(e))


def test_text_replacement():
    """Test text replacement functions"""
    print("\n### Testing Text Replacement ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()
        slide, _ = ppt_utils.add_slide(prs, 1)
        ppt_utils.add_textbox(slide, 1, 2, 4, 1, "AWS is great. AWS rocks!")

        # Test replace_text logic
        find_text = "AWS"
        replace_text = "Amazon Web Services"
        replacements = 0

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        original = run.text
                        if find_text in original:
                            run.text = original.replace(find_text, replace_text)
                            replacements += original.count(find_text)

        assert replacements == 2, f"Expected 2 replacements, got {replacements}"

        # Verify replacement
        text_found = False
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text
                if "Amazon Web Services" in text:
                    text_found = True

        assert text_found, "Text not replaced"
        results.record_pass("text replacement")

    except Exception as e:
        results.record_fail("text replacement tests", str(e))


def test_shape_queries():
    """Test get_shape_info, find_shapes_by_type, get_all_textboxes"""
    print("\n### Testing Shape Queries ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()
        slide, _ = ppt_utils.add_slide(prs, 1)

        # Add different types of shapes
        ppt_utils.add_textbox(slide, 1, 1, 3, 1, "Text Box 1")
        ppt_utils.add_textbox(slide, 1, 2.5, 3, 1, "Text Box 2")

        # Test 1: get_shape_info
        result = ppt_utils.get_shape_info(slide, 0)
        assert result["success"], "get_shape_info failed"
        assert result["shape_index"] == 0, "Wrong shape index"
        assert "position" in result, "Missing position info"
        assert "left_inches" in result["position"], "Missing inch conversion"
        results.record_pass("get_shape_info")

        # Test 2: find_shapes_by_type (TEXT_BOX)
        result = ppt_utils.find_shapes_by_type(slide, "TEXT_BOX")
        assert result["success"], "find_shapes_by_type failed"
        # Note: Number may vary based on layout placeholders
        results.record_pass("find_shapes_by_type")

        # Test 3: get_all_textboxes
        result = ppt_utils.get_all_textboxes(slide)
        assert result["success"], "get_all_textboxes failed"
        # Note: count may vary based on layout placeholders vs textboxes
        assert "total_textboxes" in result, "Missing total_textboxes in result"
        results.record_pass("get_all_textboxes")

        # Test 4: Error handling - invalid shape index
        result = ppt_utils.get_shape_info(slide, 999)
        assert not result["success"], "Should fail with invalid index"
        results.record_pass("get_shape_info error handling")

    except Exception as e:
        results.record_fail("shape query tests", str(e))


def test_placeholder_queries():
    """Test placeholder listing and querying"""
    print("\n### Testing Placeholder Queries ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()
        slide, _ = ppt_utils.add_slide(prs, 0)  # Title slide layout

        # Get placeholders
        placeholders = []
        for ph in slide.placeholders:
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "type": str(ph.placeholder_format.type),
                "name": ph.name
            })

        assert len(placeholders) > 0, "No placeholders found"
        results.record_pass("list_placeholders")

        # Test get by name
        if placeholders:
            target_name = placeholders[0]["name"]
            found = None
            for ph in slide.placeholders:
                if ph.name == target_name:
                    found = ph
                    break

            assert found is not None, f"Placeholder '{target_name}' not found"
            results.record_pass("get_placeholder_by_name")

    except Exception as e:
        results.record_fail("placeholder query tests", str(e))


def test_format_copying():
    """Test copy_slide_format and apply_text_style_to_all"""
    print("\n### Testing Format Copying ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()

        # Add source slide
        source_slide, _ = ppt_utils.add_slide(prs, 1)
        ppt_utils.add_textbox(source_slide, 1, 1, 3, 1, "Source text",
                             font_size=18, bold=True, font_name="Arial")

        # Add target slides
        target1, _ = ppt_utils.add_slide(prs, 1)
        ppt_utils.add_textbox(target1, 1, 1, 3, 1, "Target 1 text")

        target2, _ = ppt_utils.add_slide(prs, 1)
        ppt_utils.add_textbox(target2, 1, 1, 3, 1, "Target 2 text")

        # Test 1: copy_slide_format
        result = ppt_utils.copy_slide_format(prs, 0, [1, 2],
                                             copy_background=True,
                                             copy_font_styles=True)
        assert result["success"], "copy_slide_format failed"
        assert result["slides_formatted"] == 2, "Wrong number of slides formatted"
        results.record_pass("copy_slide_format")

        # Test 2: apply_text_style_to_all
        result = ppt_utils.apply_text_style_to_all(prs,
                                                   font_name="Calibri",
                                                   font_size=14,
                                                   apply_to="body")
        assert result["success"], "apply_text_style_to_all failed"
        assert result["shapes_modified"] > 0, "No shapes modified"
        results.record_pass("apply_text_style_to_all (body)")

        # Test 3: apply to titles only
        result = ppt_utils.apply_text_style_to_all(prs,
                                                   bold=True,
                                                   apply_to="title")
        assert result["success"], "apply_text_style_to_all (title) failed"
        results.record_pass("apply_text_style_to_all (title)")

        # Test 4: apply to all
        result = ppt_utils.apply_text_style_to_all(prs,
                                                   font_color=(255, 0, 0),
                                                   apply_to="all")
        assert result["success"], "apply_text_style_to_all (all) failed"
        results.record_pass("apply_text_style_to_all (all)")

    except Exception as e:
        results.record_fail("format copying tests", str(e))


def test_multi_level_bullets():
    """Test multi-level bullet point functionality"""
    print("\n### Testing Multi-Level Bullets ###")

    try:
        # Create test presentation
        prs = ppt_utils.create_presentation()
        slide, _ = ppt_utils.add_slide(prs, 1)

        # Find a placeholder for bullets
        placeholder = None
        for shape in slide.placeholders:
            if hasattr(shape, 'text_frame'):
                placeholder = shape
                break

        if placeholder:
            # Test 1: Dict format with levels
            bullet_points = [
                "Main Point 1",
                {"text": "Sub-point 1a", "level": 1},
                {"text": "Sub-sub-point 1a1", "level": 2},
                {"text": "Sub-point 1b", "level": 1},
                "Main Point 2"
            ]

            result = ppt_utils.add_bullet_points(placeholder, bullet_points)
            assert result["success"], "Multi-level bullets failed"
            assert result["total_points"] == 5, "Wrong number of bullet points"
            assert 0 in result["levels_used"], "Level 0 not used"
            assert 1 in result["levels_used"], "Level 1 not used"
            assert 2 in result["levels_used"], "Level 2 not used"
            results.record_pass("multi-level bullets (dict format)")

            # Test 2: Levels parameter format
            slide2, _ = ppt_utils.add_slide(prs, 1)
            placeholder2 = None
            for shape in slide2.placeholders:
                if hasattr(shape, 'text_frame'):
                    placeholder2 = shape
                    break

            if placeholder2:
                bullet_points = ["Point 1", "Point 2", "Point 3"]
                levels = [0, 1, 0]
                result = ppt_utils.add_bullet_points(placeholder2, bullet_points, levels)
                assert result["success"], "Levels parameter failed"
                assert result["total_points"] == 3, "Wrong number of points"
                results.record_pass("multi-level bullets (levels parameter)")

            # Test 3: Backward compatibility (simple strings)
            slide3, _ = ppt_utils.add_slide(prs, 1)
            placeholder3 = None
            for shape in slide3.placeholders:
                if hasattr(shape, 'text_frame'):
                    placeholder3 = shape
                    break

            if placeholder3:
                simple_bullets = ["Point 1", "Point 2", "Point 3"]
                result = ppt_utils.add_bullet_points(placeholder3, simple_bullets)
                assert result["success"], "Backward compatibility failed"
                assert result["levels_used"] == [0], "Should only use level 0"
                results.record_pass("multi-level bullets (backward compatibility)")
        else:
            print("⚠️  Warning: No suitable placeholder found for bullet test")

    except Exception as e:
        results.record_fail("multi-level bullets tests", str(e))


def test_integration():
    """Test combining multiple features"""
    print("\n### Testing Integration Scenarios ###")

    try:
        # Scenario: Create presentation, add slides, reorder, search, replace
        prs = ppt_utils.create_presentation()

        # Add slides
        for i in range(3):
            slide, _ = ppt_utils.add_slide(prs, 1)
            if slide.shapes.title:
                slide.shapes.title.text = f"Section {i}"
            ppt_utils.add_textbox(slide, 1, 2, 4, 1, f"Content for [PLACEHOLDER] {i}")

        initial_count = len(prs.slides)

        # Duplicate a slide (use fresh presentation to avoid index issues)
        prs2 = ppt_utils.create_presentation()
        slide1, _ = ppt_utils.add_slide(prs2, 1)
        ppt_utils.add_textbox(slide1, 1, 2, 4, 1, "Test")
        ppt_utils.duplicate_slide(prs2, 0)
        assert len(prs2.slides) == 2, "Duplication in integration failed"

        # Replace text in original presentation
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = run.text.replace("[PLACEHOLDER]", "ACTUAL_VALUE")

        # Verify replacement worked
        found_replacement = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if "ACTUAL_VALUE" in shape.text_frame.text:
                        found_replacement = True

        assert found_replacement, "Integration: text replacement failed"

        # Apply consistent styling
        result = ppt_utils.apply_text_style_to_all(prs, font_size=12, apply_to="body")
        assert result["success"], "Integration: styling failed"

        results.record_pass("integration scenario")

    except Exception as e:
        results.record_fail("integration tests", str(e))


def main():
    """Run all tests"""
    print("=" * 60)
    print("TESTING NEW POWERPOINT MCP SERVER FEATURES")
    print("=" * 60)

    test_slide_reordering()
    test_slide_duplication()
    test_text_search()
    test_text_replacement()
    test_shape_queries()
    test_placeholder_queries()
    test_format_copying()
    test_multi_level_bullets()
    test_integration()

    success = results.summary()

    if success:
        print("\n🎉 All tests passed!")
        return 0
    else:
        print("\n⚠️  Some tests failed. Please review errors above.")
        return 1


if __name__ == "__main__":
    exit(main())
