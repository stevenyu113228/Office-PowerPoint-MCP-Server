"""
Shape alignment tools for PowerPoint MCP Server.
Handles shape alignment, distribution, and layout analysis.
"""
from typing import Dict, List, Optional, Any, Tuple
from mcp.server.fastmcp import FastMCP


def register_shape_alignment_tools(app: FastMCP, presentations: Dict, get_current_presentation_id, validate_parameters, is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register shape alignment tools with the FastMCP app"""
    
    # ---- Batch Alignment Functions ----
    
    @app.tool()
    def align_shapes(
        slide_index: int,
        shape_indices: List[int],
        alignment_type: str,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Align multiple shapes using various alignment types."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        # Validate shape indices
        for shape_index in shape_indices:
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                }
        
        if len(shape_indices) < 2:
            return {
                "error": "At least 2 shapes are required for alignment"
            }
        
        valid_alignments = ["top", "bottom", "left", "right", "center_horizontal", "center_vertical", "center"]
        if alignment_type not in valid_alignments:
            return {
                "error": f"Invalid alignment type: '{alignment_type}'. Valid types: {', '.join(valid_alignments)}"
            }
        
        try:
            from pptx.util import Inches
            shapes = [slide.shapes[i] for i in shape_indices]
            
            # Store original positions
            original_positions = []
            for i, shape in enumerate(shapes):
                original_positions.append({
                    "index": shape_indices[i],
                    "left": shape.left / 914400,
                    "top": shape.top / 914400
                })
            
            if alignment_type == "top":
                # Align to the topmost shape
                min_top = min(shape.top for shape in shapes)
                for shape in shapes:
                    shape.top = min_top
                    
            elif alignment_type == "bottom":
                # Align to the bottom of the bottommost shape
                max_bottom = max(shape.top + shape.height for shape in shapes)
                for shape in shapes:
                    shape.top = max_bottom - shape.height
                    
            elif alignment_type == "left":
                # Align to the leftmost shape
                min_left = min(shape.left for shape in shapes)
                for shape in shapes:
                    shape.left = min_left
                    
            elif alignment_type == "right":
                # Align to the right of the rightmost shape
                max_right = max(shape.left + shape.width for shape in shapes)
                for shape in shapes:
                    shape.left = max_right - shape.width
                    
            elif alignment_type == "center_horizontal":
                # Align to horizontal center
                if len(shapes) > 0:
                    center_x = sum(shape.left + shape.width / 2 for shape in shapes) / len(shapes)
                    for shape in shapes:
                        shape.left = int(center_x - shape.width / 2)
                    
            elif alignment_type == "center_vertical":
                # Align to vertical center
                if len(shapes) > 0:
                    center_y = sum(shape.top + shape.height / 2 for shape in shapes) / len(shapes)
                    for shape in shapes:
                        shape.top = int(center_y - shape.height / 2)
                    
            elif alignment_type == "center":
                # Align to both horizontal and vertical center
                if len(shapes) > 0:
                    center_x = sum(shape.left + shape.width / 2 for shape in shapes) / len(shapes)
                    center_y = sum(shape.top + shape.height / 2 for shape in shapes) / len(shapes)
                for shape in shapes:
                    shape.left = int(center_x - shape.width / 2)
                    shape.top = int(center_y - shape.height / 2)
            
            # Get new positions
            new_positions = []
            for i, shape in enumerate(shapes):
                new_positions.append({
                    "index": shape_indices[i],
                    "left": shape.left / 914400,
                    "top": shape.top / 914400
                })
            
            return {
                "message": f"Aligned {len(shapes)} shapes using '{alignment_type}' alignment on slide {slide_index}",
                "alignment_type": alignment_type,
                "shape_indices": shape_indices,
                "original_positions": original_positions,
                "new_positions": new_positions
            }
        except Exception as e:
            return {
                "error": f"Failed to align shapes: {str(e)}"
            }
    
    @app.tool()
    def align_shapes_horizontally(
        slide_index: int,
        shape_indices: List[int],
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Align shapes horizontally (center alignment)."""
        return align_shapes(slide_index, shape_indices, "center_horizontal", presentation_id)
    
    @app.tool()
    def align_shapes_vertically(
        slide_index: int,
        shape_indices: List[int],
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Align shapes vertically (center alignment)."""
        return align_shapes(slide_index, shape_indices, "center_vertical", presentation_id)
    
    @app.tool()
    def distribute_shapes_evenly(
        slide_index: int,
        shape_indices: List[int],
        direction: str = "horizontal",
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Distribute shapes evenly in horizontal or vertical direction."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        # Validate shape indices
        for shape_index in shape_indices:
            if shape_index < 0 or shape_index >= len(slide.shapes):
                return {
                    "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
                }
        
        if len(shape_indices) < 3:
            return {
                "error": "At least 3 shapes are required for distribution"
            }
        
        if direction not in ["horizontal", "vertical"]:
            return {
                "error": f"Invalid direction: '{direction}'. Valid directions: 'horizontal', 'vertical'"
            }
        
        try:
            from pptx.util import Inches
            shapes = [slide.shapes[i] for i in shape_indices]
            
            # Store original positions
            original_positions = []
            for i, shape in enumerate(shapes):
                original_positions.append({
                    "index": shape_indices[i],
                    "left": shape.left / 914400,
                    "top": shape.top / 914400
                })
            
            if direction == "horizontal":
                # Sort shapes by left position
                sorted_shapes = sorted(zip(shapes, shape_indices), key=lambda x: x[0].left)
                
                # Calculate total space and gaps
                leftmost = sorted_shapes[0][0].left
                rightmost = sorted_shapes[-1][0].left + sorted_shapes[-1][0].width
                total_width = rightmost - leftmost
                
                # Calculate total shape widths
                total_shape_width = sum(shape.width for shape, _ in sorted_shapes)
                
                # Calculate gap between shapes
                if len(sorted_shapes) > 1:
                    gap = (total_width - total_shape_width) / (len(sorted_shapes) - 1)
                else:
                    gap = 0
                
                # Distribute shapes
                current_x = leftmost
                for shape, _ in sorted_shapes:
                    shape.left = int(current_x)
                    current_x += shape.width + gap
                    
            else:  # vertical
                # Sort shapes by top position
                sorted_shapes = sorted(zip(shapes, shape_indices), key=lambda x: x[0].top)
                
                # Calculate total space and gaps
                topmost = sorted_shapes[0][0].top
                bottommost = sorted_shapes[-1][0].top + sorted_shapes[-1][0].height
                total_height = bottommost - topmost
                
                # Calculate total shape heights
                total_shape_height = sum(shape.height for shape, _ in sorted_shapes)
                
                # Calculate gap between shapes
                if len(sorted_shapes) > 1:
                    gap = (total_height - total_shape_height) / (len(sorted_shapes) - 1)
                else:
                    gap = 0
                
                # Distribute shapes
                current_y = topmost
                for shape, _ in sorted_shapes:
                    shape.top = int(current_y)
                    current_y += shape.height + gap
            
            # Get new positions
            new_positions = []
            for i, shape in enumerate(shapes):
                new_positions.append({
                    "index": shape_indices[i],
                    "left": shape.left / 914400,
                    "top": shape.top / 914400
                })
            
            return {
                "message": f"Distributed {len(shapes)} shapes evenly in {direction} direction on slide {slide_index}",
                "direction": direction,
                "shape_indices": shape_indices,
                "original_positions": original_positions,
                "new_positions": new_positions
            }
        except Exception as e:
            return {
                "error": f"Failed to distribute shapes: {str(e)}"
            }
    
    # ---- Layout Analysis Functions ----
    
    @app.tool()
    def detect_overlapping_elements(
        slide_index: int,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Detect overlapping elements on a slide."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            overlaps = []
            shapes = list(slide.shapes)
            
            for i in range(len(shapes)):
                for j in range(i + 1, len(shapes)):
                    shape1 = shapes[i]
                    shape2 = shapes[j]
                    
                    # Calculate bounds
                    s1_left = shape1.left / 914400
                    s1_top = shape1.top / 914400
                    s1_right = s1_left + (shape1.width / 914400)
                    s1_bottom = s1_top + (shape1.height / 914400)
                    
                    s2_left = shape2.left / 914400
                    s2_top = shape2.top / 914400
                    s2_right = s2_left + (shape2.width / 914400)
                    s2_bottom = s2_top + (shape2.height / 914400)
                    
                    # Check for overlap
                    if not (s1_right <= s2_left or s1_left >= s2_right or 
                           s1_bottom <= s2_top or s1_top >= s2_bottom):
                        
                        # Calculate overlap area
                        overlap_left = max(s1_left, s2_left)
                        overlap_top = max(s1_top, s2_top)
                        overlap_right = min(s1_right, s2_right)
                        overlap_bottom = min(s1_bottom, s2_bottom)
                        
                        overlap_width = overlap_right - overlap_left
                        overlap_height = overlap_bottom - overlap_top
                        overlap_area = overlap_width * overlap_height
                        
                        overlaps.append({
                            "shape1_index": i,
                            "shape1_name": getattr(shape1, 'name', f"Shape {i}"),
                            "shape2_index": j,
                            "shape2_name": getattr(shape2, 'name', f"Shape {j}"),
                            "overlap_area": overlap_area,
                            "overlap_bounds": {
                                "left": overlap_left,
                                "top": overlap_top,
                                "width": overlap_width,
                                "height": overlap_height
                            }
                        })
            
            return {
                "slide_index": slide_index,
                "total_shapes": len(shapes),
                "overlapping_pairs": len(overlaps),
                "overlaps": overlaps,
                "has_overlaps": len(overlaps) > 0
            }
        except Exception as e:
            return {
                "error": f"Failed to detect overlapping elements: {str(e)}"
            }
    
    @app.tool()
    def suggest_alignment_improvements(
        slide_index: int,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Suggest alignment improvements for a slide."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            suggestions = []
            shapes = list(slide.shapes)
            
            if len(shapes) < 2:
                return {
                    "slide_index": slide_index,
                    "suggestions": [],
                    "message": "Not enough shapes to suggest alignment improvements"
                }
            
            # Analyze horizontal alignment
            left_positions = [shape.left / 914400 for shape in shapes]
            right_positions = [(shape.left + shape.width) / 914400 for shape in shapes]
            center_x_positions = [(shape.left + shape.width / 2) / 914400 for shape in shapes]
            
            # Check for near-alignments (within 0.1 inches)
            tolerance = 0.1
            
            # Check left alignment opportunities
            for i in range(len(shapes)):
                near_left = [j for j in range(len(shapes)) if j != i and abs(left_positions[i] - left_positions[j]) <= tolerance]
                if len(near_left) >= 1:
                    indices = [i] + near_left
                    suggestions.append({
                        "type": "left_alignment",
                        "shape_indices": indices,
                        "description": f"Shapes {indices} could be left-aligned (currently within {tolerance} inches)"
                    })
            
            # Check right alignment opportunities
            for i in range(len(shapes)):
                near_right = [j for j in range(len(shapes)) if j != i and abs(right_positions[i] - right_positions[j]) <= tolerance]
                if len(near_right) >= 1:
                    indices = [i] + near_right
                    suggestions.append({
                        "type": "right_alignment",
                        "shape_indices": indices,
                        "description": f"Shapes {indices} could be right-aligned (currently within {tolerance} inches)"
                    })
            
            # Check center alignment opportunities
            for i in range(len(shapes)):
                near_center = [j for j in range(len(shapes)) if j != i and abs(center_x_positions[i] - center_x_positions[j]) <= tolerance]
                if len(near_center) >= 1:
                    indices = [i] + near_center
                    suggestions.append({
                        "type": "center_horizontal_alignment",
                        "shape_indices": indices,
                        "description": f"Shapes {indices} could be center-aligned horizontally (currently within {tolerance} inches)"
                    })
            
            # Analyze vertical alignment
            top_positions = [shape.top / 914400 for shape in shapes]
            bottom_positions = [(shape.top + shape.height) / 914400 for shape in shapes]
            center_y_positions = [(shape.top + shape.height / 2) / 914400 for shape in shapes]
            
            # Check top alignment opportunities
            for i in range(len(shapes)):
                near_top = [j for j in range(len(shapes)) if j != i and abs(top_positions[i] - top_positions[j]) <= tolerance]
                if len(near_top) >= 1:
                    indices = [i] + near_top
                    suggestions.append({
                        "type": "top_alignment",
                        "shape_indices": indices,
                        "description": f"Shapes {indices} could be top-aligned (currently within {tolerance} inches)"
                    })
            
            # Check for distribution opportunities
            if len(shapes) >= 3:
                # Check if shapes are roughly in a line and could be distributed
                horizontal_shapes = sorted(enumerate(shapes), key=lambda x: x[1].left)
                if len(horizontal_shapes) >= 3:
                    # Check if they're roughly horizontally aligned
                    top_variance = max(top_positions) - min(top_positions)
                    if top_variance <= tolerance * 2:
                        suggestions.append({
                            "type": "horizontal_distribution",
                            "shape_indices": [i for i, _ in horizontal_shapes],
                            "description": "Shapes could be distributed evenly horizontally"
                        })
                
                vertical_shapes = sorted(enumerate(shapes), key=lambda x: x[1].top)
                if len(vertical_shapes) >= 3:
                    # Check if they're roughly vertically aligned
                    left_variance = max(left_positions) - min(left_positions)
                    if left_variance <= tolerance * 2:
                        suggestions.append({
                            "type": "vertical_distribution",
                            "shape_indices": [i for i, _ in vertical_shapes],
                            "description": "Shapes could be distributed evenly vertically"
                        })
            
            # Remove duplicate suggestions
            unique_suggestions = []
            seen = set()
            for suggestion in suggestions:
                key = (suggestion["type"], tuple(sorted(suggestion["shape_indices"])))
                if key not in seen:
                    seen.add(key)
                    unique_suggestions.append(suggestion)
            
            return {
                "slide_index": slide_index,
                "total_shapes": len(shapes),
                "suggestions": unique_suggestions,
                "suggestion_count": len(unique_suggestions),
                "tolerance": tolerance
            }
        except Exception as e:
            return {
                "error": f"Failed to suggest alignment improvements: {str(e)}"
            }
    
    @app.tool()
    def get_layout_metrics(
        slide_index: int,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Get layout metrics for a slide including spacing and alignment analysis."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            shapes = list(slide.shapes)
            
            if len(shapes) == 0:
                return {
                    "slide_index": slide_index,
                    "total_shapes": 0,
                    "message": "No shapes on slide to analyze"
                }
            
            # Get slide dimensions
            slide_width = pres.slide_width / 914400
            slide_height = pres.slide_height / 914400
            
            # Calculate shape positions and sizes
            shape_metrics = []
            for i, shape in enumerate(shapes):
                left = shape.left / 914400
                top = shape.top / 914400
                width = shape.width / 914400
                height = shape.height / 914400
                
                shape_metrics.append({
                    "index": i,
                    "name": getattr(shape, 'name', f"Shape {i}"),
                    "left": left,
                    "top": top,
                    "right": left + width,
                    "bottom": top + height,
                    "width": width,
                    "height": height,
                    "center_x": left + width / 2,
                    "center_y": top + height / 2,
                    "area": width * height
                })
            
            # Calculate spacing metrics
            horizontal_gaps = []
            vertical_gaps = []
            
            for i in range(len(shape_metrics)):
                for j in range(i + 1, len(shape_metrics)):
                    shape1 = shape_metrics[i]
                    shape2 = shape_metrics[j]
                    
                    # Horizontal gap (if shapes are side by side)
                    if (shape1["top"] < shape2["bottom"] and shape1["bottom"] > shape2["top"]):
                        if shape1["right"] <= shape2["left"]:
                            horizontal_gaps.append(shape2["left"] - shape1["right"])
                        elif shape2["right"] <= shape1["left"]:
                            horizontal_gaps.append(shape1["left"] - shape2["right"])
                    
                    # Vertical gap (if shapes are above/below each other)
                    if (shape1["left"] < shape2["right"] and shape1["right"] > shape2["left"]):
                        if shape1["bottom"] <= shape2["top"]:
                            vertical_gaps.append(shape2["top"] - shape1["bottom"])
                        elif shape2["bottom"] <= shape1["top"]:
                            vertical_gaps.append(shape1["top"] - shape2["bottom"])
            
            # Calculate alignment metrics
            left_positions = [s["left"] for s in shape_metrics]
            right_positions = [s["right"] for s in shape_metrics]
            top_positions = [s["top"] for s in shape_metrics]
            bottom_positions = [s["bottom"] for s in shape_metrics]
            center_x_positions = [s["center_x"] for s in shape_metrics]
            center_y_positions = [s["center_y"] for s in shape_metrics]
            
            # Calculate coverage
            if shape_metrics:
                content_left = min(left_positions)
                content_right = max(right_positions)
                content_top = min(top_positions)
                content_bottom = max(bottom_positions)
                
                content_width = content_right - content_left
                content_height = content_bottom - content_top
                content_area = sum(s["area"] for s in shape_metrics)
                
                coverage_x = (content_width / slide_width) * 100
                coverage_y = (content_height / slide_height) * 100
                density = (content_area / (slide_width * slide_height)) * 100
            else:
                coverage_x = coverage_y = density = 0
                content_left = content_right = content_top = content_bottom = 0
                content_width = content_height = 0
            
            return {
                "slide_index": slide_index,
                "slide_dimensions": {
                    "width": slide_width,
                    "height": slide_height,
                    "aspect_ratio": slide_width / slide_height if slide_height != 0 else 0
                },
                "total_shapes": len(shapes),
                "shape_metrics": shape_metrics,
                "spacing_analysis": {
                    "horizontal_gaps": horizontal_gaps,
                    "vertical_gaps": vertical_gaps,
                    "avg_horizontal_gap": sum(horizontal_gaps) / len(horizontal_gaps) if horizontal_gaps else 0,
                    "avg_vertical_gap": sum(vertical_gaps) / len(vertical_gaps) if vertical_gaps else 0,
                    "min_horizontal_gap": min(horizontal_gaps) if horizontal_gaps else 0,
                    "max_horizontal_gap": max(horizontal_gaps) if horizontal_gaps else 0,
                    "min_vertical_gap": min(vertical_gaps) if vertical_gaps else 0,
                    "max_vertical_gap": max(vertical_gaps) if vertical_gaps else 0
                },
                "alignment_analysis": {
                    "left_positions_variance": max(left_positions) - min(left_positions) if left_positions else 0,
                    "right_positions_variance": max(right_positions) - min(right_positions) if right_positions else 0,
                    "top_positions_variance": max(top_positions) - min(top_positions) if top_positions else 0,
                    "bottom_positions_variance": max(bottom_positions) - min(bottom_positions) if bottom_positions else 0,
                    "center_x_variance": max(center_x_positions) - min(center_x_positions) if center_x_positions else 0,
                    "center_y_variance": max(center_y_positions) - min(center_y_positions) if center_y_positions else 0
                },
                "coverage_analysis": {
                    "content_bounds": {
                        "left": content_left,
                        "top": content_top,
                        "right": content_right,
                        "bottom": content_bottom,
                        "width": content_width,
                        "height": content_height
                    },
                    "coverage_x_percent": coverage_x,
                    "coverage_y_percent": coverage_y,
                    "density_percent": density
                }
            }
        except Exception as e:
            return {
                "error": f"Failed to get layout metrics: {str(e)}"
            }