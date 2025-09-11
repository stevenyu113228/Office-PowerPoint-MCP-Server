"""
Shape positioning and manipulation tools for PowerPoint MCP Server.
Handles shape position, size, alignment, and coordinate system operations.
"""
from typing import Dict, List, Optional, Any, Tuple
from mcp.server.fastmcp import FastMCP


def register_shape_positioning_tools(app: FastMCP, presentations: Dict, get_current_presentation_id, validate_parameters, is_positive, is_non_negative, is_in_range, is_valid_rgb):
    """Register shape positioning and manipulation tools with the FastMCP app"""
    
    # ---- Coordinate System Utilities ----
    
    @app.tool()
    def convert_emu_to_inches(emu_value: int) -> Dict:
        """Convert EMU (English Metric Units) to inches."""
        try:
            # Validate input
            if not isinstance(emu_value, (int, float)):
                return {
                    "error": "EMU value must be a number"
                }
            
            # EMU to inches conversion: 1 inch = 914400 EMU
            inches = float(emu_value) / 914400
            return {
                "emu": emu_value,
                "inches": inches,
                "conversion_factor": 914400
            }
        except Exception as e:
            return {
                "error": f"Failed to convert EMU to inches: {str(e)}"
            }
    
    @app.tool()
    def convert_inches_to_emu(inch_value: float) -> Dict:
        """Convert inches to EMU (English Metric Units)."""
        try:
            # Validate input
            if not isinstance(inch_value, (int, float)):
                return {
                    "error": "Inch value must be a number"
                }
            
            # Inches to EMU conversion: 1 inch = 914400 EMU
            emu = int(float(inch_value) * 914400)
            return {
                "inches": inch_value,
                "emu": emu,
                "conversion_factor": 914400
            }
        except Exception as e:
            return {
                "error": f"Failed to convert inches to EMU: {str(e)}"
            }
    
    @app.tool()
    def get_slide_dimensions_inches(slide_index: int, presentation_id: Optional[str] = None) -> Dict:
        """Get slide dimensions in inches."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        try:
            # Get slide dimensions from presentation
            slide_width_emu = pres.slide_width
            slide_height_emu = pres.slide_height
            
            # Convert to inches
            slide_width_inches = slide_width_emu / 914400
            slide_height_inches = slide_height_emu / 914400
            
            # Calculate aspect ratio with zero division protection
            aspect_ratio = slide_width_inches / slide_height_inches if slide_height_inches != 0 else 0
            
            return {
                "slide_index": slide_index,
                "width_inches": slide_width_inches,
                "height_inches": slide_height_inches,
                "width_emu": slide_width_emu,
                "height_emu": slide_height_emu,
                "aspect_ratio": aspect_ratio
            }
        except Exception as e:
            return {
                "error": f"Failed to get slide dimensions: {str(e)}"
            }
    
    # ---- Shape Position and Size Adjustment ----
    
    @app.tool()
    def update_shape_position(
        slide_index: int,
        shape_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Update shape position and size."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        # Validate parameters
        validations = {
            "left": (left, [(is_non_negative, "must be non-negative")]),
            "top": (top, [(is_non_negative, "must be non-negative")]),
            "width": (width, [(is_positive, "must be positive")]),
            "height": (height, [(is_positive, "must be positive")])
        }
        
        valid, error = validate_parameters(validations)
        if not valid:
            return {"error": error}
        
        try:
            from pptx.util import Inches
            shape = slide.shapes[shape_index]
            
            # Store original position and size
            original = {
                "left": shape.left / 914400,
                "top": shape.top / 914400,
                "width": shape.width / 914400,
                "height": shape.height / 914400
            }
            
            # Update position and size
            shape.left = Inches(left)
            shape.top = Inches(top)
            shape.width = Inches(width)
            shape.height = Inches(height)
            
            return {
                "message": f"Updated shape {shape_index} position and size on slide {slide_index}",
                "shape_index": shape_index,
                "original": original,
                "new": {
                    "left": left,
                    "top": top,
                    "width": width,
                    "height": height
                }
            }
        except Exception as e:
            return {
                "error": f"Failed to update shape position: {str(e)}"
            }
    
    @app.tool()
    def update_shape_size(
        slide_index: int,
        shape_index: int,
        width: float,
        height: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Update shape size only."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        # Validate parameters
        validations = {
            "width": (width, [(is_positive, "must be positive")]),
            "height": (height, [(is_positive, "must be positive")])
        }
        
        valid, error = validate_parameters(validations)
        if not valid:
            return {"error": error}
        
        try:
            from pptx.util import Inches
            shape = slide.shapes[shape_index]
            
            # Store original size
            original_width = shape.width / 914400
            original_height = shape.height / 914400
            
            # Update size
            shape.width = Inches(width)
            shape.height = Inches(height)
            
            return {
                "message": f"Updated shape {shape_index} size on slide {slide_index}",
                "shape_index": shape_index,
                "original_size": {
                    "width": original_width,
                    "height": original_height
                },
                "new_size": {
                    "width": width,
                    "height": height
                }
            }
        except Exception as e:
            return {
                "error": f"Failed to update shape size: {str(e)}"
            }
    
    @app.tool()
    def move_shape(
        slide_index: int,
        shape_index: int,
        delta_x: float,
        delta_y: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Move shape by specified delta values."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        try:
            from pptx.util import Inches
            shape = slide.shapes[shape_index]
            
            # Store original position
            original_left = shape.left / 914400
            original_top = shape.top / 914400
            
            # Calculate new position
            new_left = original_left + delta_x
            new_top = original_top + delta_y
            
            # Ensure new position is not negative
            new_left = max(0, new_left)
            new_top = max(0, new_top)
            
            # Update position
            shape.left = Inches(new_left)
            shape.top = Inches(new_top)
            
            return {
                "message": f"Moved shape {shape_index} on slide {slide_index}",
                "shape_index": shape_index,
                "delta": {
                    "x": delta_x,
                    "y": delta_y
                },
                "original_position": {
                    "left": original_left,
                    "top": original_top
                },
                "new_position": {
                    "left": new_left,
                    "top": new_top
                }
            }
        except Exception as e:
            return {
                "error": f"Failed to move shape: {str(e)}"
            }
    
    # ---- Image/Picture Operations ----
    
    @app.tool()
    def update_image_position(
        slide_index: int,
        shape_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Update image/picture position and size."""
        # This is the same as update_shape_position but specifically for images
        return update_shape_position(slide_index, shape_index, left, top, width, height, presentation_id)
    
    @app.tool()
    def resize_image(
        slide_index: int,
        shape_index: int,
        width: float,
        height: float,
        maintain_aspect_ratio: bool = True,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Resize image with optional aspect ratio maintenance."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        try:
            from pptx.util import Inches
            shape = slide.shapes[shape_index]
            
            # Check if it's an image
            if not hasattr(shape, 'image'):
                return {
                    "error": f"Shape at index {shape_index} is not an image"
                }
            
            # Store original dimensions
            original_width = shape.width / 914400
            original_height = shape.height / 914400
            original_aspect_ratio = original_width / original_height if original_height != 0 else 1
            
            if maintain_aspect_ratio:
                # Calculate new dimensions maintaining aspect ratio
                new_aspect_ratio = width / height if height != 0 else 1
                
                if new_aspect_ratio > original_aspect_ratio:
                    # Width is the limiting factor
                    final_height = height
                    final_width = height * original_aspect_ratio
                else:
                    # Height is the limiting factor
                    final_width = width
                    final_height = width / original_aspect_ratio if original_aspect_ratio != 0 else width
            else:
                final_width = width
                final_height = height
            
            # Update size
            shape.width = Inches(final_width)
            shape.height = Inches(final_height)
            
            return {
                "message": f"Resized image {shape_index} on slide {slide_index}",
                "shape_index": shape_index,
                "maintain_aspect_ratio": maintain_aspect_ratio,
                "original_size": {
                    "width": original_width,
                    "height": original_height,
                    "aspect_ratio": original_aspect_ratio
                },
                "requested_size": {
                    "width": width,
                    "height": height
                },
                "final_size": {
                    "width": final_width,
                    "height": final_height,
                    "aspect_ratio": final_width / final_height if final_height != 0 else 0
                }
            }
        except Exception as e:
            return {
                "error": f"Failed to resize image: {str(e)}"
            }
    
    # ---- TextBox Operations ----
    
    @app.tool()
    def update_textbox_position(
        slide_index: int,
        shape_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Update textbox position and size."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        try:
            shape = slide.shapes[shape_index]
            
            # Check if it has text frame (textbox)
            if not hasattr(shape, 'text_frame'):
                return {
                    "error": f"Shape at index {shape_index} is not a textbox"
                }
            
            # Use the existing update_shape_position function
            return update_shape_position(slide_index, shape_index, left, top, width, height, presentation_id)
            
        except Exception as e:
            return {
                "error": f"Failed to update textbox position: {str(e)}"
            }
    
    @app.tool()
    def move_textbox(
        slide_index: int,
        shape_index: int,
        delta_x: float,
        delta_y: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Move textbox by specified delta values."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if shape_index < 0 or shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        try:
            shape = slide.shapes[shape_index]
            
            # Check if it has text frame (textbox)
            if not hasattr(shape, 'text_frame'):
                return {
                    "error": f"Shape at index {shape_index} is not a textbox"
                }
            
            # Use the existing move_shape function
            return move_shape(slide_index, shape_index, delta_x, delta_y, presentation_id)
            
        except Exception as e:
            return {
                "error": f"Failed to move textbox: {str(e)}"
            }
    
    # ---- AutoShape Operations ----
    
    @app.tool()
    def update_autoshape_position(
        slide_index: int,
        shape_index: int,
        left: float,
        top: float,
        width: float,
        height: float,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Update autoshape position and size."""
        # This is the same as update_shape_position but specifically for autoshapes
        return update_shape_position(slide_index, shape_index, left, top, width, height, presentation_id)
    
    # ---- Shape Query and Selection ----
    
    @app.tool()
    def get_shapes_by_type(
        slide_index: int,
        shape_type: str,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Get shapes by type."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            from pptx.shapes.autoshape import AutoShape
            from pptx.shapes.picture import Picture
            from pptx.shapes.table import Table
            from pptx.shapes.graphfrm import GraphicFrame
            
            matching_shapes = []
            shape_type_lower = shape_type.lower()
            
            for i, shape in enumerate(slide.shapes):
                shape_info = {
                    "index": i,
                    "left": shape.left / 914400,
                    "top": shape.top / 914400,
                    "width": shape.width / 914400,
                    "height": shape.height / 914400
                }
                
                # Determine shape type
                if shape_type_lower == "autoshape" and isinstance(shape, AutoShape):
                    shape_info["shape_type"] = "autoshape"
                    if hasattr(shape, 'auto_shape_type'):
                        shape_info["auto_shape_type"] = str(shape.auto_shape_type)
                    matching_shapes.append(shape_info)
                elif shape_type_lower == "picture" and isinstance(shape, Picture):
                    shape_info["shape_type"] = "picture"
                    matching_shapes.append(shape_info)
                elif shape_type_lower == "table" and isinstance(shape, Table):
                    shape_info["shape_type"] = "table"
                    shape_info["rows"] = len(shape.table.rows)
                    shape_info["columns"] = len(shape.table.columns)
                    matching_shapes.append(shape_info)
                elif shape_type_lower == "chart" and isinstance(shape, GraphicFrame):
                    if hasattr(shape, 'chart'):
                        shape_info["shape_type"] = "chart"
                        shape_info["chart_type"] = str(shape.chart.chart_type)
                        matching_shapes.append(shape_info)
                elif shape_type_lower == "textbox" and hasattr(shape, 'text_frame'):
                    shape_info["shape_type"] = "textbox"
                    shape_info["text"] = shape.text_frame.text[:100] + "..." if len(shape.text_frame.text) > 100 else shape.text_frame.text
                    matching_shapes.append(shape_info)
            
            return {
                "slide_index": slide_index,
                "shape_type": shape_type,
                "matching_shapes": matching_shapes,
                "count": len(matching_shapes)
            }
        except Exception as e:
            return {
                "error": f"Failed to get shapes by type: {str(e)}"
            }
    
    @app.tool()
    def get_shapes_by_name_pattern(
        slide_index: int,
        pattern: str,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Get shapes by name pattern."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        try:
            import re
            matching_shapes = []
            
            for i, shape in enumerate(slide.shapes):
                shape_name = getattr(shape, 'name', f"Shape {i}")
                
                # Check if pattern matches the shape name
                if re.search(pattern, shape_name, re.IGNORECASE):
                    shape_info = {
                        "index": i,
                        "name": shape_name,
                        "left": shape.left / 914400,
                        "top": shape.top / 914400,
                        "width": shape.width / 914400,
                        "height": shape.height / 914400
                    }
                    
                    # Add text content if it's a text shape
                    if hasattr(shape, 'text_frame') and shape.text_frame.text:
                        shape_info["text"] = shape.text_frame.text[:100] + "..." if len(shape.text_frame.text) > 100 else shape.text_frame.text
                    
                    matching_shapes.append(shape_info)
            
            return {
                "slide_index": slide_index,
                "pattern": pattern,
                "matching_shapes": matching_shapes,
                "count": len(matching_shapes)
            }
        except Exception as e:
            return {
                "error": f"Failed to get shapes by name pattern: {str(e)}"
            }
    
    @app.tool()
    def get_overlapping_shapes(
        slide_index: int,
        reference_shape_index: int,
        presentation_id: Optional[str] = None
    ) -> Dict:
        """Find shapes that overlap with a reference shape."""
        pres_id = presentation_id if presentation_id is not None else get_current_presentation_id()
        
        if pres_id is None or pres_id not in presentations:
            return {
                "error": "No presentation is currently loaded or the specified ID is invalid"
            }
        
        pres = presentations[pres_id]
        
        if slide_index < 0 or slide_index >= len(pres.slides):
            return {
                "error": f"Invalid slide index: {slide_index}. Available slides: 0-{len(pres.slides) - 1}"
            }
        
        slide = pres.slides[slide_index]
        
        if reference_shape_index < 0 or reference_shape_index >= len(slide.shapes):
            return {
                "error": f"Invalid reference shape index: {reference_shape_index}. Available shapes: 0-{len(slide.shapes) - 1}"
            }
        
        try:
            ref_shape = slide.shapes[reference_shape_index]
            ref_left = ref_shape.left / 914400
            ref_top = ref_shape.top / 914400
            ref_right = ref_left + (ref_shape.width / 914400)
            ref_bottom = ref_top + (ref_shape.height / 914400)
            
            overlapping_shapes = []
            
            for i, shape in enumerate(slide.shapes):
                if i == reference_shape_index:
                    continue  # Skip the reference shape itself
                
                shape_left = shape.left / 914400
                shape_top = shape.top / 914400
                shape_right = shape_left + (shape.width / 914400)
                shape_bottom = shape_top + (shape.height / 914400)
                
                # Check for overlap
                if not (shape_right <= ref_left or shape_left >= ref_right or 
                       shape_bottom <= ref_top or shape_top >= ref_bottom):
                    
                    # Calculate overlap area
                    overlap_left = max(ref_left, shape_left)
                    overlap_top = max(ref_top, shape_top)
                    overlap_right = min(ref_right, shape_right)
                    overlap_bottom = min(ref_bottom, shape_bottom)
                    
                    overlap_width = overlap_right - overlap_left
                    overlap_height = overlap_bottom - overlap_top
                    overlap_area = overlap_width * overlap_height
                    
                    shape_info = {
                        "index": i,
                        "name": getattr(shape, 'name', f"Shape {i}"),
                        "left": shape_left,
                        "top": shape_top,
                        "width": shape.width / 914400,
                        "height": shape.height / 914400,
                        "overlap_area": overlap_area,
                        "overlap_percentage": (overlap_area / ((ref_shape.width / 914400) * (ref_shape.height / 914400))) * 100
                    }
                    overlapping_shapes.append(shape_info)
            
            return {
                "slide_index": slide_index,
                "reference_shape_index": reference_shape_index,
                "reference_shape": {
                    "left": ref_left,
                    "top": ref_top,
                    "width": ref_shape.width / 914400,
                    "height": ref_shape.height / 914400
                },
                "overlapping_shapes": overlapping_shapes,
                "count": len(overlapping_shapes)
            }
        except Exception as e:
            return {
                "error": f"Failed to find overlapping shapes: {str(e)}"
            }