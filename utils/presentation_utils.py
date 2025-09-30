"""
Presentation management utilities for PowerPoint MCP Server.
Functions for creating, opening, saving, and managing presentations.
"""
from pptx import Presentation
from pptx.util import Inches
from typing import Dict, List, Optional
import os
import copy
from io import BytesIO


def create_presentation() -> Presentation:
    """
    Create a new PowerPoint presentation.
    
    Returns:
        A new Presentation object
    """
    return Presentation()


def open_presentation(file_path: str) -> Presentation:
    """
    Open an existing PowerPoint presentation.
    
    Args:
        file_path: Path to the PowerPoint file
        
    Returns:
        A Presentation object
    """
    return Presentation(file_path)


def create_presentation_from_template(template_path: str) -> Presentation:
    """
    Create a new PowerPoint presentation from a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        A new Presentation object based on the template
        
    Raises:
        FileNotFoundError: If the template file doesn't exist
        Exception: If the template file is corrupted or invalid
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    if not template_path.lower().endswith(('.pptx', '.potx')):
        raise ValueError("Template file must be a .pptx or .potx file")
    
    try:
        # Load the template file as a presentation
        presentation = Presentation(template_path)
        return presentation
    except Exception as e:
        raise Exception(f"Failed to load template file '{template_path}': {str(e)}")


def save_presentation(presentation: Presentation, file_path: str) -> str:
    """
    Save a PowerPoint presentation to a file.
    
    Args:
        presentation: The Presentation object
        file_path: Path where the file should be saved
        
    Returns:
        The file path where the presentation was saved
    """
    presentation.save(file_path)
    return file_path


def get_template_info(template_path: str) -> Dict:
    """
    Get information about a template file.
    
    Args:
        template_path: Path to the template .pptx file
        
    Returns:
        Dictionary containing template information
    """
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template file not found: {template_path}")
    
    try:
        presentation = Presentation(template_path)
        
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        # Get file size
        file_size = os.path.getsize(template_path)
        
        return {
            "template_path": template_path,
            "file_size_bytes": file_size,
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props
        }
    except Exception as e:
        raise Exception(f"Failed to read template info from '{template_path}': {str(e)}")


def get_presentation_info(presentation: Presentation) -> Dict:
    """
    Get information about a presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        Dictionary containing presentation information
    """
    try:
        # Get slide layouts
        layouts = get_slide_layouts(presentation)
        
        # Get core properties
        core_props = get_core_properties(presentation)
        
        # Get slide count
        slide_count = len(presentation.slides)
        
        return {
            "slide_count": slide_count,
            "layout_count": len(layouts),
            "slide_layouts": layouts,
            "core_properties": core_props,
            "slide_width": presentation.slide_width,
            "slide_height": presentation.slide_height
        }
    except Exception as e:
        raise Exception(f"Failed to get presentation info: {str(e)}")


def get_slide_layouts(presentation: Presentation) -> List[Dict]:
    """
    Get all available slide layouts in the presentation.
    
    Args:
        presentation: The Presentation object
        
    Returns:
        A list of dictionaries with layout information
    """
    layouts = []
    for i, layout in enumerate(presentation.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        }
        layouts.append(layout_info)
    return layouts


def set_core_properties(presentation: Presentation, title: str = None, subject: str = None,
                       author: str = None, keywords: str = None, comments: str = None) -> None:
    """
    Set core document properties.
    
    Args:
        presentation: The Presentation object
        title: Document title
        subject: Document subject
        author: Document author
        keywords: Document keywords
        comments: Document comments
    """
    core_props = presentation.core_properties
    
    if title is not None:
        core_props.title = title
    if subject is not None:
        core_props.subject = subject
    if author is not None:
        core_props.author = author
    if keywords is not None:
        core_props.keywords = keywords
    if comments is not None:
        core_props.comments = comments


def get_core_properties(presentation: Presentation) -> Dict:
    """
    Get core document properties.

    Args:
        presentation: The Presentation object

    Returns:
        Dictionary containing core properties
    """
    core_props = presentation.core_properties

    return {
        "title": core_props.title,
        "subject": core_props.subject,
        "author": core_props.author,
        "keywords": core_props.keywords,
        "comments": core_props.comments,
        "created": core_props.created.isoformat() if core_props.created else None,
        "last_modified_by": core_props.last_modified_by,
        "modified": core_props.modified.isoformat() if core_props.modified else None
    }


# Slide Reordering Functions

def move_slide(presentation: Presentation, from_index: int, to_index: int) -> Dict:
    """
    Move a slide from one position to another.

    Args:
        presentation: The Presentation object
        from_index: Current index of the slide (0-based)
        to_index: Target index for the slide (0-based)

    Returns:
        Dictionary with operation results
    """
    slides = presentation.slides
    total_slides = len(slides)

    # Validate indices
    if from_index < 0 or from_index >= total_slides:
        raise ValueError(f"Invalid from_index: {from_index}. Must be 0-{total_slides - 1}")
    if to_index < 0 or to_index >= total_slides:
        raise ValueError(f"Invalid to_index: {to_index}. Must be 0-{total_slides - 1}")

    if from_index == to_index:
        return {
            "success": True,
            "message": "No move needed - indices are the same",
            "from_index": from_index,
            "to_index": to_index
        }

    # Get the XML element for the slide
    slide = slides[from_index]
    xml_slides = presentation.slides._sldIdLst

    # Move the slide in the XML
    # Store reference to the slide element
    slide_elem = xml_slides[from_index]

    # Remove from original position
    xml_slides.remove(slide_elem)

    # Insert at new position
    xml_slides.insert(to_index, slide_elem)

    return {
        "success": True,
        "message": f"Moved slide from position {from_index} to {to_index}",
        "from_index": from_index,
        "to_index": to_index
    }


def swap_slides(presentation: Presentation, index_a: int, index_b: int) -> Dict:
    """
    Swap two slides.

    Args:
        presentation: The Presentation object
        index_a: Index of first slide
        index_b: Index of second slide

    Returns:
        Dictionary with operation results
    """
    total_slides = len(presentation.slides)

    # Validate indices
    if index_a < 0 or index_a >= total_slides:
        raise ValueError(f"Invalid index_a: {index_a}. Must be 0-{total_slides - 1}")
    if index_b < 0 or index_b >= total_slides:
        raise ValueError(f"Invalid index_b: {index_b}. Must be 0-{total_slides - 1}")

    if index_a == index_b:
        return {
            "success": True,
            "message": "No swap needed - indices are the same",
            "index_a": index_a,
            "index_b": index_b
        }

    # Use simple approach: move one slide at a time
    if index_a > index_b:
        # Move higher index first to avoid shifting issues
        move_slide(presentation, index_a, index_b)
        move_slide(presentation, index_b + 1, index_a)
    else:
        # Move lower index first
        move_slide(presentation, index_b, index_a)
        move_slide(presentation, index_a + 1, index_b)

    return {
        "success": True,
        "message": f"Swapped slides at positions {index_a} and {index_b}",
        "index_a": index_a,
        "index_b": index_b
    }


def reorder_slides(presentation: Presentation, new_order: List[int]) -> Dict:
    """
    Reorder slides according to a new arrangement.

    Args:
        presentation: The Presentation object
        new_order: List of slide indices in desired order
                  Example: [0, 3, 1, 2, 4] means keep slide 0 first,
                          move slide 3 to position 1, etc.

    Returns:
        Dictionary with operation results
    """
    total_slides = len(presentation.slides)

    # Validate new_order
    if len(new_order) != total_slides:
        raise ValueError(f"new_order must contain {total_slides} indices, got {len(new_order)}")

    if sorted(new_order) != list(range(total_slides)):
        raise ValueError(f"new_order must contain all indices 0-{total_slides - 1} exactly once")

    xml_slides = presentation.slides._sldIdLst

    # Create a copy of the current order
    original_order = list(xml_slides)

    # Clear the current order
    for _ in range(len(xml_slides)):
        xml_slides.remove(xml_slides[0])

    # Add slides in new order
    for idx in new_order:
        xml_slides.append(original_order[idx])

    return {
        "success": True,
        "message": f"Reordered {total_slides} slides",
        "new_order": new_order,
        "total_slides": total_slides
    }


# Slide Deletion Functions

def delete_slide(presentation: Presentation, slide_index: int) -> Dict:
    """
    Delete a slide from the presentation.

    Args:
        presentation: The Presentation object
        slide_index: Index of the slide to delete (0-based)

    Returns:
        Dictionary with operation results
    """
    total_slides = len(presentation.slides)

    # Validate index
    if slide_index < 0 or slide_index >= total_slides:
        raise ValueError(f"Invalid slide_index: {slide_index}. Must be 0-{total_slides - 1}")

    # Get the XML slides list
    xml_slides = presentation.slides._sldIdLst
    rId = xml_slides[slide_index].rId

    # Remove the slide from the XML
    xml_slides.remove(xml_slides[slide_index])

    return {
        "success": True,
        "message": f"Deleted slide at index {slide_index}",
        "deleted_index": slide_index,
        "remaining_slides": len(presentation.slides)
    }


def delete_slides(presentation: Presentation, slide_indices: List[int]) -> Dict:
    """
    Delete multiple slides from the presentation.

    Args:
        presentation: The Presentation object
        slide_indices: List of slide indices to delete (0-based)

    Returns:
        Dictionary with operation results
    """
    total_slides = len(presentation.slides)

    # Validate all indices first
    for idx in slide_indices:
        if idx < 0 or idx >= total_slides:
            raise ValueError(f"Invalid slide index: {idx}. Must be 0-{total_slides - 1}")

    # Check for duplicates
    if len(slide_indices) != len(set(slide_indices)):
        raise ValueError("Duplicate indices found in slide_indices")

    # Sort indices in descending order to delete from end to start
    # This prevents index shifting issues
    sorted_indices = sorted(slide_indices, reverse=True)

    deleted_count = 0
    for idx in sorted_indices:
        try:
            delete_slide(presentation, idx)
            deleted_count += 1
        except Exception as e:
            # Continue deleting other slides even if one fails
            pass

    return {
        "success": True,
        "message": f"Deleted {deleted_count} slides",
        "deleted_count": deleted_count,
        "requested_count": len(slide_indices),
        "remaining_slides": len(presentation.slides)
    }


# Slide Duplication Functions

def duplicate_slide(presentation: Presentation, slide_index: int,
                   insert_position: Optional[int] = None) -> Dict:
    """
    Duplicate a slide within the presentation.

    Args:
        presentation: The Presentation object
        slide_index: Index of the slide to duplicate
        insert_position: Where to insert the duplicate (None = after original)

    Returns:
        Dictionary with operation results including new slide index
    """
    total_slides = len(presentation.slides)

    # Validate slide_index
    if slide_index < 0 or slide_index >= total_slides:
        raise ValueError(f"Invalid slide_index: {slide_index}. Must be 0-{total_slides - 1}")

    # Default insert position
    if insert_position is None:
        insert_position = slide_index + 1
    elif insert_position < 0 or insert_position > total_slides:
        raise ValueError(f"Invalid insert_position: {insert_position}. Must be 0-{total_slides}")

    source_slide = presentation.slides[slide_index]

    # Get the slide layout
    slide_layout = source_slide.slide_layout

    # Create a new slide with the same layout
    blank_slide = presentation.slides.add_slide(slide_layout)

    # Copy all shapes from source to new slide
    for shape in source_slide.shapes:
        _copy_shape_to_slide(shape, blank_slide)

    # Move the new slide to the correct position if needed
    new_slide_index = len(presentation.slides) - 1
    if insert_position != new_slide_index:
        move_slide(presentation, new_slide_index, insert_position)
        new_slide_index = insert_position

    return {
        "success": True,
        "message": f"Duplicated slide {slide_index} to position {new_slide_index}",
        "original_index": slide_index,
        "new_index": new_slide_index
    }


def _copy_shape_to_slide(source_shape, target_slide):
    """
    Helper function to copy a shape to another slide.
    This is a simplified version - some complex shapes may not copy perfectly.
    """
    try:
        # Copy text boxes and shapes with text
        if source_shape.has_text_frame:
            # Create similar shape
            new_shape = target_slide.shapes.add_textbox(
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )

            # Copy text frame content
            source_tf = source_shape.text_frame
            target_tf = new_shape.text_frame

            target_tf.text = source_tf.text
            target_tf.word_wrap = source_tf.word_wrap

            # Copy paragraph formatting
            for src_para, tgt_para in zip(source_tf.paragraphs, target_tf.paragraphs):
                tgt_para.alignment = src_para.alignment
                tgt_para.level = src_para.level

                # Copy run formatting
                for src_run, tgt_run in zip(src_para.runs, tgt_para.runs):
                    tgt_run.font.bold = src_run.font.bold
                    tgt_run.font.italic = src_run.font.italic
                    if src_run.font.size:
                        tgt_run.font.size = src_run.font.size
                    if src_run.font.name:
                        tgt_run.font.name = src_run.font.name

        # Copy pictures
        elif source_shape.shape_type == 13:  # PICTURE
            # Pictures need to be handled via their image part
            image_part = source_shape.image.blob
            target_slide.shapes.add_picture(
                BytesIO(image_part),
                source_shape.left,
                source_shape.top,
                source_shape.width,
                source_shape.height
            )
    except Exception:
        # Skip shapes that can't be copied
        pass