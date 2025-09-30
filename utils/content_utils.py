"""
Content management utilities for PowerPoint MCP Server.
Functions for slides, text, images, tables, charts, and shapes.
"""
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from typing import Dict, List, Tuple, Optional, Any
import tempfile
import os
import base64


def add_slide(presentation: Presentation, layout_index: int = 1) -> Tuple:
    """
    Add a slide to the presentation.
    
    Args:
        presentation: The Presentation object
        layout_index: Index of the slide layout to use
        
    Returns:
        A tuple containing the slide and its layout
    """
    layout = presentation.slide_layouts[layout_index]
    slide = presentation.slides.add_slide(layout)
    return slide, layout


def get_slide_info(slide, slide_index: int) -> Dict:
    """
    Get information about a specific slide.
    
    Args:
        slide: The slide object
        slide_index: Index of the slide
        
    Returns:
        Dictionary containing slide information
    """
    try:
        placeholders = []
        for placeholder in slide.placeholders:
            placeholder_info = {
                "idx": placeholder.placeholder_format.idx,
                "type": str(placeholder.placeholder_format.type),
                "name": placeholder.name
            }
            placeholders.append(placeholder_info)
        
        shapes = []
        for i, shape in enumerate(slide.shapes):
            shape_info = {
                "index": i,
                "name": shape.name,
                "shape_type": str(shape.shape_type),
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height
            }
            shapes.append(shape_info)
        
        return {
            "slide_index": slide_index,
            "layout_name": slide.slide_layout.name,
            "placeholder_count": len(placeholders),
            "placeholders": placeholders,
            "shape_count": len(shapes),
            "shapes": shapes
        }
    except Exception as e:
        raise Exception(f"Failed to get slide info: {str(e)}")


def set_title(slide, title: str) -> None:
    """
    Set the title of a slide.
    
    Args:
        slide: The slide object
        title: The title text
    """
    if slide.shapes.title:
        slide.shapes.title.text = title


def populate_placeholder(slide, placeholder_idx: int, text: str) -> None:
    """
    Populate a placeholder with text.
    
    Args:
        slide: The slide object
        placeholder_idx: The index of the placeholder
        text: The text to add
    """
    placeholder = slide.placeholders[placeholder_idx]
    placeholder.text = text


def add_bullet_points(placeholder, bullet_points: List[Any], levels: Optional[List[int]] = None) -> Dict:
    """
    Add bullet points to a placeholder with optional multi-level indentation.

    Supports both simple strings and dict format with indentation levels:
    - Simple format: List of strings (all level 0, backward compatible)
    - Advanced format: List of strings or dicts with "text" and "level" keys
    - Separate levels: List of strings with separate levels parameter

    Args:
        placeholder: The placeholder object
        bullet_points: List of bullet point texts or dicts with {"text": str, "level": int}
        levels: Optional list of indentation levels (0-8) if bullet_points is List[str]

    Returns:
        Dict with operation results including levels used

    Example:
        # Simple usage (backward compatible)
        add_bullet_points(placeholder, ["Point 1", "Point 2"])

        # With indentation using dict format
        add_bullet_points(placeholder, [
            "Main Point",
            {"text": "Sub-point", "level": 1},
            {"text": "Sub-sub-point", "level": 2}
        ])

        # With indentation using separate levels parameter
        add_bullet_points(placeholder, ["Point 1", "Sub 1", "Sub 2"], levels=[0, 1, 1])
    """
    text_frame = placeholder.text_frame
    text_frame.clear()

    levels_used = set()

    for i, point in enumerate(bullet_points):
        # Determine text and level
        if isinstance(point, dict):
            text = point.get("text", "")
            level = point.get("level", 0)
        else:
            text = point
            level = levels[i] if levels and i < len(levels) else 0

        # Ensure level is within valid range (0-8 per PowerPoint specification)
        level = max(0, min(8, level))
        levels_used.add(level)

        # Add paragraph
        if i == 0:
            text_frame.text = text
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
            p.text = text

        # Set indentation level (core feature)
        p.level = level

        # Optional: Adjust font size based on level for better visual hierarchy
        if level == 0:
            if p.runs:
                p.runs[0].font.size = Pt(18)
        elif level == 1:
            if p.runs:
                p.runs[0].font.size = Pt(16)
        elif level >= 2:
            if p.runs:
                p.runs[0].font.size = Pt(14)

    return {
        "success": True,
        "total_points": len(bullet_points),
        "levels_used": sorted(list(levels_used))
    }


def add_textbox(slide, left: float, top: float, width: float, height: float, text: str,
                font_size: int = None, font_name: str = None, bold: bool = None,
                italic: bool = None, underline: bool = None, 
                color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                alignment: str = None, vertical_alignment: str = None, 
                auto_fit: bool = True) -> Any:
    """
    Add a textbox to a slide with formatting options.
    
    Args:
        slide: The slide object
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        text: Text content
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
        auto_fit: Whether to auto-fit text
        
    Returns:
        The created textbox shape
    """
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    
    textbox.text_frame.text = text
    
    # Apply formatting if provided
    if any([font_size, font_name, bold, italic, underline, color, bg_color, alignment, vertical_alignment]):
        format_text_advanced(
            textbox.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            underline=underline,
            color=color,
            bg_color=bg_color,
            alignment=alignment,
            vertical_alignment=vertical_alignment
        )
    
    return textbox


def format_text(text_frame, font_size: int = None, font_name: str = None, 
                bold: bool = None, italic: bool = None, color: Tuple[int, int, int] = None,
                alignment: str = None) -> None:
    """
    Format text in a text frame.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
    """
    alignment_map = {
        'left': PP_ALIGN.LEFT,
        'center': PP_ALIGN.CENTER,
        'right': PP_ALIGN.RIGHT,
        'justify': PP_ALIGN.JUSTIFY
    }
    
    for paragraph in text_frame.paragraphs:
        if alignment and alignment in alignment_map:
            paragraph.alignment = alignment_map[alignment]
            
        for run in paragraph.runs:
            font = run.font
            
            if font_size is not None:
                font.size = Pt(font_size)
            if font_name is not None:
                font.name = font_name
            if bold is not None:
                font.bold = bold
            if italic is not None:
                font.italic = italic
            if color is not None:
                r, g, b = color
                font.color.rgb = RGBColor(r, g, b)


def format_text_advanced(text_frame, font_size: int = None, font_name: str = None, 
                        bold: bool = None, italic: bool = None, underline: bool = None,
                        color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                        alignment: str = None, vertical_alignment: str = None) -> Dict:
    """
    Advanced text formatting with comprehensive options.
    
    Args:
        text_frame: The text frame to format
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        underline: Whether text should be underlined
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment ('left', 'center', 'right', 'justify')
        vertical_alignment: Vertical alignment ('top', 'middle', 'bottom')
    
    Returns:
        Dictionary with formatting results
    """
    result = {
        'success': True,
        'warnings': []
    }
    
    try:
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        
        # Enable text wrapping
        text_frame.word_wrap = True
        
        # Apply formatting to all paragraphs and runs
        for paragraph in text_frame.paragraphs:
            if alignment and alignment in alignment_map:
                paragraph.alignment = alignment_map[alignment]
            
            for run in paragraph.runs:
                font = run.font
                
                if font_size is not None:
                    font.size = Pt(font_size)
                if font_name is not None:
                    font.name = font_name
                if bold is not None:
                    font.bold = bold
                if italic is not None:
                    font.italic = italic
                if underline is not None:
                    font.underline = underline
                if color is not None:
                    r, g, b = color
                    font.color.rgb = RGBColor(r, g, b)
        
        return result
        
    except Exception as e:
        result['success'] = False
        result['error'] = str(e)
        return result


def add_image(slide, image_path: str, left: float, top: float, width: float = None, height: float = None) -> Any:
    """
    Add an image to a slide.
    
    Args:
        slide: The slide object
        image_path: Path to the image file
        left: Left position in inches
        top: Top position in inches
        width: Width in inches (optional)
        height: Height in inches (optional)
        
    Returns:
        The created image shape
    """
    if width is not None and height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width), Inches(height)
        )
    elif width is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), Inches(width)
        )
    elif height is not None:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top), height=Inches(height)
        )
    else:
        return slide.shapes.add_picture(
            image_path, Inches(left), Inches(top)
        )


def add_table(slide, rows: int, cols: int, left: float, top: float, width: float, height: float) -> Any:
    """
    Add a table to a slide.
    
    Args:
        slide: The slide object
        rows: Number of rows
        cols: Number of columns
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        
    Returns:
        The created table shape
    """
    return slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
    )


def format_table_cell(cell, font_size: int = None, font_name: str = None, 
                     bold: bool = None, italic: bool = None, 
                     color: Tuple[int, int, int] = None, bg_color: Tuple[int, int, int] = None,
                     alignment: str = None, vertical_alignment: str = None) -> None:
    """
    Format a table cell.
    
    Args:
        cell: The table cell object
        font_size: Font size in points
        font_name: Font name
        bold: Whether text should be bold
        italic: Whether text should be italic
        color: RGB color tuple (r, g, b)
        bg_color: Background RGB color tuple (r, g, b)
        alignment: Text alignment
        vertical_alignment: Vertical alignment
    """
    # Format text
    if any([font_size, font_name, bold, italic, color, alignment]):
        format_text_advanced(
            cell.text_frame,
            font_size=font_size,
            font_name=font_name,
            bold=bold,
            italic=italic,
            color=color,
            alignment=alignment
        )
    
    # Set background color
    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*bg_color)


def add_chart(slide, chart_type: str, left: float, top: float, width: float, height: float,
              categories: List[str], series_names: List[str], series_values: List[List[float]]) -> Any:
    """
    Add a chart to a slide.
    
    Args:
        slide: The slide object
        chart_type: Type of chart ('column', 'bar', 'line', 'pie', etc.)
        left: Left position in inches
        top: Top position in inches
        width: Width in inches
        height: Height in inches
        categories: List of category names
        series_names: List of series names
        series_values: List of value lists for each series
        
    Returns:
        The created chart object
    """
    # Map chart type names to enum values
    chart_type_map = {
        'column': XL_CHART_TYPE.COLUMN_CLUSTERED,
        'stacked_column': XL_CHART_TYPE.COLUMN_STACKED,
        'bar': XL_CHART_TYPE.BAR_CLUSTERED,
        'stacked_bar': XL_CHART_TYPE.BAR_STACKED,
        'line': XL_CHART_TYPE.LINE,
        'line_markers': XL_CHART_TYPE.LINE_MARKERS,
        'pie': XL_CHART_TYPE.PIE,
        'doughnut': XL_CHART_TYPE.DOUGHNUT,
        'area': XL_CHART_TYPE.AREA,
        'stacked_area': XL_CHART_TYPE.AREA_STACKED,
        'scatter': XL_CHART_TYPE.XY_SCATTER,
        'radar': XL_CHART_TYPE.RADAR,
        'radar_markers': XL_CHART_TYPE.RADAR_MARKERS
    }
    
    xl_chart_type = chart_type_map.get(chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = categories
    
    for i, series_name in enumerate(series_names):
        if i < len(series_values):
            chart_data.add_series(series_name, series_values[i])
    
    # Add chart to slide
    chart_shape = slide.shapes.add_chart(
        xl_chart_type, Inches(left), Inches(top), Inches(width), Inches(height), chart_data
    )
    
    return chart_shape.chart


def format_chart(chart, has_legend: bool = True, legend_position: str = 'right',
                has_data_labels: bool = False, title: str = None,
                x_axis_title: str = None, y_axis_title: str = None,
                color_scheme: str = None) -> None:
    """
    Format a chart with various options.
    
    Args:
        chart: The chart object
        has_legend: Whether to show legend
        legend_position: Position of legend ('right', 'top', 'bottom', 'left')
        has_data_labels: Whether to show data labels
        title: Chart title
        x_axis_title: X-axis title
        y_axis_title: Y-axis title
        color_scheme: Color scheme to apply
    """
    try:
        # Set chart title
        if title:
            chart.chart_title.text_frame.text = title
        
        # Configure legend
        if has_legend:
            chart.has_legend = True
            # Note: Legend position setting may vary by chart type
        else:
            chart.has_legend = False
        
        # Configure data labels
        if has_data_labels:
            for series in chart.series:
                series.has_data_labels = True
        
        # Set axis titles if available
        try:
            if x_axis_title and hasattr(chart, 'category_axis'):
                chart.category_axis.axis_title.text_frame.text = x_axis_title
            if y_axis_title and hasattr(chart, 'value_axis'):
                chart.value_axis.axis_title.text_frame.text = y_axis_title
        except:
            pass  # Axis titles may not be available for all chart types
            
    except Exception:
        pass  # Graceful degradation for chart formatting


def format_keywords_in_text(shape, keywords: List[str],
                           bold: Optional[bool] = None,
                           italic: Optional[bool] = None,
                           underline: Optional[bool] = None,
                           font_size: Optional[int] = None,
                           font_color: Optional[Tuple[int, int, int]] = None,
                           case_sensitive: bool = False) -> Dict:
    """
    Format specific keywords in a text shape with different styling.

    Args:
        shape: The shape object containing text
        keywords: List of keywords to format
        bold: Whether to make keywords bold
        italic: Whether to make keywords italic
        underline: Whether to underline keywords
        font_size: Font size for keywords in points
        font_color: RGB color tuple for keywords
        case_sensitive: Whether search is case-sensitive

    Returns:
        Dictionary with operation results
    """
    try:
        if not hasattr(shape, 'text_frame') or not shape.text_frame:
            return {
                "success": False,
                "error": "Shape does not have a text frame"
            }

        text_frame = shape.text_frame
        total_formatted = 0
        keywords_found = {}

        # Process each paragraph
        for paragraph in text_frame.paragraphs:
            # Get full paragraph text
            para_text = paragraph.text

            # Find keyword positions in paragraph
            for keyword in keywords:
                keyword_count = 0
                search_text = para_text if case_sensitive else para_text.lower()
                search_keyword = keyword if case_sensitive else keyword.lower()

                # Find all occurrences
                start = 0
                positions = []
                while True:
                    pos = search_text.find(search_keyword, start)
                    if pos == -1:
                        break
                    positions.append((pos, pos + len(keyword)))
                    start = pos + 1
                    keyword_count += 1

                if keyword_count > 0:
                    keywords_found[keyword] = keyword_count
                    total_formatted += keyword_count

                # Apply formatting to found keywords
                if positions:
                    # Need to rebuild paragraph with runs
                    _apply_formatting_to_positions(paragraph, positions, para_text,
                                                  bold, italic, underline,
                                                  font_size, font_color)

        return {
            "success": True,
            "total_keywords_formatted": total_formatted,
            "keywords_found": keywords_found,
            "message": f"Formatted {total_formatted} keyword occurrences"
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to format keywords: {str(e)}"
        }


def _apply_formatting_to_positions(paragraph, positions, original_text,
                                   bold, italic, underline, font_size, font_color):
    """
    Helper function to apply formatting to specific positions in a paragraph.
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    # Clear existing runs
    for _ in range(len(paragraph.runs)):
        paragraph._p.remove(paragraph.runs[0]._r)

    # Sort positions
    positions = sorted(positions)

    # Build new runs with formatting
    last_end = 0
    for start, end in positions:
        # Add unformatted text before keyword
        if start > last_end:
            run = paragraph.add_run()
            run.text = original_text[last_end:start]

        # Add formatted keyword
        run = paragraph.add_run()
        run.text = original_text[start:end]

        # Apply formatting
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic
        if underline is not None:
            run.font.underline = underline
        if font_size is not None:
            run.font.size = Pt(font_size)
        if font_color is not None:
            run.font.color.rgb = RGBColor(*font_color)

        last_end = end

    # Add remaining unformatted text
    if last_end < len(original_text):
        run = paragraph.add_run()
        run.text = original_text[last_end:]


def get_shape_info(slide, shape_index: int) -> Dict:
    """
    Get detailed information about a specific shape.

    Args:
        slide: The slide object
        shape_index: Index of the shape

    Returns:
        Dictionary with detailed shape information
    """
    try:
        if shape_index < 0 or shape_index >= len(slide.shapes):
            raise ValueError(f"Invalid shape index: {shape_index}. Available shapes: 0-{len(slide.shapes) - 1}")

        shape = slide.shapes[shape_index]

        info = {
            "shape_index": shape_index,
            "name": shape.name,
            "type": str(shape.shape_type),
            "position": {
                "left": shape.left,
                "top": shape.top,
                "width": shape.width,
                "height": shape.height,
                "left_inches": shape.left / 914400,
                "top_inches": shape.top / 914400,
                "width_inches": shape.width / 914400,
                "height_inches": shape.height / 914400
            },
            "has_text": hasattr(shape, 'text_frame') and shape.text_frame is not None,
            "is_placeholder": hasattr(shape, 'placeholder_format')
        }

        # Add text information if available
        if info["has_text"]:
            text = shape.text_frame.text
            info["text_length"] = len(text)
            info["text_preview"] = text[:100] + "..." if len(text) > 100 else text

            # Get font information from first run if available
            if shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
                first_run = shape.text_frame.paragraphs[0].runs[0]
                font_info = {}
                if first_run.font.name:
                    font_info["name"] = first_run.font.name
                if first_run.font.size:
                    font_info["size"] = int(first_run.font.size / 12700)  # Convert to points
                if first_run.font.bold is not None:
                    font_info["bold"] = first_run.font.bold
                if first_run.font.italic is not None:
                    font_info["italic"] = first_run.font.italic
                info["font"] = font_info

        # Add placeholder information if applicable
        if info["is_placeholder"]:
            info["placeholder"] = {
                "idx": shape.placeholder_format.idx,
                "type": str(shape.placeholder_format.type)
            }

        return {
            "success": True,
            **info
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to get shape info: {str(e)}"
        }


def find_shapes_by_type(slide, shape_type: str) -> Dict:
    """
    Find all shapes of a specific type on a slide.

    Args:
        slide: The slide object
        shape_type: Type of shape to find (e.g., "TEXT_BOX", "PICTURE", "TABLE", "CHART")

    Returns:
        Dictionary with matching shapes
    """
    try:
        matches = []

        for i, shape in enumerate(slide.shapes):
            # Match by shape type name
            if shape_type.upper() in str(shape.shape_type).upper():
                matches.append({
                    "shape_index": i,
                    "name": shape.name,
                    "type": str(shape.shape_type),
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                })

        return {
            "success": True,
            "shape_type": shape_type,
            "total_matches": len(matches),
            "matches": matches
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to find shapes: {str(e)}"
        }


def get_all_textboxes(slide) -> Dict:
    """
    Get all text boxes from a slide.

    Args:
        slide: The slide object

    Returns:
        Dictionary with all text boxes and their content
    """
    try:
        textboxes = []

        for i, shape in enumerate(slide.shapes):
            # Check if it's a textbox (has text_frame but is not a placeholder)
            is_placeholder = False
            try:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    is_placeholder = True
            except:
                # Some shapes throw errors when accessing placeholder_format
                pass

            if hasattr(shape, 'text_frame') and shape.text_frame and not is_placeholder:
                text = shape.text_frame.text
                textboxes.append({
                    "shape_index": i,
                    "name": shape.name,
                    "text": text,
                    "text_length": len(text),
                    "position": {
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height
                    }
                })

        return {
            "success": True,
            "total_textboxes": len(textboxes),
            "textboxes": textboxes
        }

    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to get textboxes: {str(e)}"
        }


def extract_slide_text_content(slide) -> Dict:
    """
    Extract all text content from a slide including placeholders and text shapes.
    
    Args:
        slide: The slide object to extract text from
        
    Returns:
        Dictionary containing all text content organized by source type
    """
    try:
        text_content = {
            "slide_title": "",
            "placeholders": [],
            "text_shapes": [],
            "table_text": [],
            "all_text_combined": ""
        }
        
        all_texts = []
        
        # Extract title from slide if available
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            try:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    text_content["slide_title"] = title_text
                    all_texts.append(title_text)
            except:
                pass
        
        # Extract text from all shapes
        for i, shape in enumerate(slide.shapes):
            shape_text_info = {
                "shape_index": i,
                "shape_name": shape.name,
                "shape_type": str(shape.shape_type),
                "text": ""
            }
            
            try:
                # Check if shape has text frame
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        shape_text_info["text"] = text
                        all_texts.append(text)
                        
                        # Categorize by shape type
                        if hasattr(shape, 'placeholder_format'):
                            # This is a placeholder
                            placeholder_info = shape_text_info.copy()
                            placeholder_info["placeholder_type"] = str(shape.placeholder_format.type)
                            placeholder_info["placeholder_idx"] = shape.placeholder_format.idx
                            text_content["placeholders"].append(placeholder_info)
                        else:
                            # This is a regular text shape
                            text_content["text_shapes"].append(shape_text_info)
                
                # Extract text from tables
                elif hasattr(shape, 'table'):
                    table_texts = []
                    table = shape.table
                    for row_idx, row in enumerate(table.rows):
                        row_texts = []
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text_frame.text.strip()
                            if cell_text:
                                row_texts.append(cell_text)
                                all_texts.append(cell_text)
                        if row_texts:
                            table_texts.append({
                                "row": row_idx,
                                "cells": row_texts
                            })
                    
                    if table_texts:
                        text_content["table_text"].append({
                            "shape_index": i,
                            "shape_name": shape.name,
                            "table_content": table_texts
                        })
                        
            except Exception as e:
                # Skip shapes that can't be processed
                continue
        
        # Combine all text
        text_content["all_text_combined"] = "\n".join(all_texts)
        
        return {
            "success": True,
            "text_content": text_content,
            "total_text_shapes": len(text_content["placeholders"]) + len(text_content["text_shapes"]),
            "has_title": bool(text_content["slide_title"]),
            "has_tables": len(text_content["table_text"]) > 0
        }
        
    except Exception as e:
        return {
            "success": False,
            "error": f"Failed to extract text content: {str(e)}",
            "text_content": None
        }